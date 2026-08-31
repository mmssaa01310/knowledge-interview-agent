from __future__ import annotations

import logging
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Any

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from ai_interviewer_api.core.config import settings
from ai_interviewer_api.models.base import utc_now
from ai_interviewer_api.repositories.document_knowledge import (
    DocumentKnowledgeBackendError,
    document_knowledge_repository,
)
from ai_interviewer_api.repositories.store import store

logger = logging.getLogger(__name__)

SUPPORTED_DOCUMENT_EXTENSIONS = frozenset({
    ".csv",
    ".docx",
    ".md",
    ".pdf",
    ".pptx",
    ".txt",
    ".xlsx",
})


class DocumentIngestionError(ValueError):
    """Raised when a document cannot be extracted or indexed."""


class UnsupportedDocumentTypeError(DocumentIngestionError):
    """Raised when a file type is not supported by the local extractor."""


@dataclass(frozen=True)
class DocumentIngestionResult:
    document: dict[str, Any]
    content: str
    chunks: list[dict[str, Any]]


def ingest_document(
    document: dict[str, Any],
    raw_bytes: bytes,
) -> DocumentIngestionResult:
    """Extract, chunk, and index one uploaded document.

    The document metadata remains in the application store. Only the
    searchable text and chunks are delegated to the configured document
    knowledge repository, so the same ingestion flow works with PostgreSQL
    and Elastic Cloud.
    """

    _update_document(document["id"], ingestionStatus="processing", progressPercent=20)
    try:
        content = extract_document_text(
            file_name=str(document.get("fileName") or ""),
            content_type=str(document.get("contentType") or ""),
            raw_bytes=raw_bytes,
        )
        _update_document(
            document["id"],
            ingestionStatus="text_extracted",
            progressPercent=50,
            errorMessage=None,
        )
        chunks = chunk_document_text(document, content)
        _update_document(
            document["id"],
            ingestionStatus="chunked",
            progressPercent=70,
            chunkCount=len(chunks),
        )
        document_knowledge_repository.replace_document(
            document,
            content=content,
            chunks=chunks,
        )
        result_document = _update_document(
            document["id"],
            ingestionStatus="indexed",
            progressPercent=100,
            chunkCount=len(chunks),
            lastIngestedAt=utc_now(),
            errorMessage=None,
        )
        return DocumentIngestionResult(
            document=result_document,
            content=content,
            chunks=chunks,
        )
    except DocumentKnowledgeBackendError:
        _update_document(
            document["id"],
            ingestionStatus="failed",
            errorMessage="document_backend_unavailable",
        )
        raise
    except DocumentIngestionError as error:
        _update_document(
            document["id"],
            ingestionStatus="failed",
            errorMessage=str(error),
        )
        raise
    except Exception as error:  # noqa: BLE001 - keep ingestion failures user-safe
        logger.exception("document_ingestion_failed document_id=%s", document["id"])
        _update_document(
            document["id"],
            ingestionStatus="failed",
            errorMessage="document_ingestion_failed",
        )
        raise DocumentIngestionError("document_ingestion_failed") from error


def extract_document_text(*, file_name: str, content_type: str, raw_bytes: bytes) -> str:
    extension = _document_extension(file_name)
    if extension not in SUPPORTED_DOCUMENT_EXTENSIONS:
        raise UnsupportedDocumentTypeError(
            f"unsupported_document_type:{extension or content_type or 'unknown'}"
        )
    if not raw_bytes:
        raise DocumentIngestionError("document_content_empty")

    try:
        if extension in {".csv", ".md", ".txt"}:
            text = _decode_text(raw_bytes)
        elif extension == ".pdf":
            text = "\n".join(page.extract_text() or "" for page in PdfReader(BytesIO(raw_bytes)).pages)
        elif extension == ".docx":
            text = _extract_docx_text(raw_bytes)
        elif extension == ".xlsx":
            text = _extract_xlsx_text(raw_bytes)
        elif extension == ".pptx":
            text = _extract_pptx_text(raw_bytes)
        else:  # pragma: no cover - guarded by SUPPORTED_DOCUMENT_EXTENSIONS
            raise UnsupportedDocumentTypeError(f"unsupported_document_type:{extension}")
    except UnsupportedDocumentTypeError:
        raise
    except Exception as error:  # noqa: BLE001 - parser-specific errors are user-safe
        raise DocumentIngestionError("document_text_extraction_failed") from error

    normalized = "\n".join(line.strip() for line in text.replace("\x00", "").splitlines()).strip()
    if not normalized:
        raise DocumentIngestionError("document_content_empty")
    return normalized


def chunk_document_text(
    document: dict[str, Any],
    content: str,
) -> list[dict[str, Any]]:
    chunk_size = max(100, int(settings.document_chunk_size_chars))
    overlap = max(0, min(int(settings.document_chunk_overlap_chars), chunk_size - 1))
    chunks: list[dict[str, Any]] = []
    start = 0
    chunk_number = 1
    while start < len(content):
        end = min(len(content), start + chunk_size)
        chunk_content = content[start:end].strip()
        if chunk_content:
            chunks.append(
                {
                    "id": f"{document['id']}:chunk:{chunk_number}",
                    "tenantId": document["tenantId"],
                    "createdByUserId": document.get("createdByUserId"),
                    "updatedByUserId": document.get("updatedByUserId"),
                    "knowledgeId": document["knowledgeId"],
                    "documentId": document["id"],
                    "title": document.get("fileName") or "事前知識チャンク",
                    "sequence": chunk_number,
                    "status": "indexed",
                    "ingestionStatus": "indexed",
                    "content": chunk_content,
                    "createdAt": utc_now(),
                    "updatedAt": utc_now(),
                    "deletedAt": None,
                }
            )
            chunk_number += 1
        if end >= len(content):
            break
        start = end - overlap
    if not chunks:
        raise DocumentIngestionError("document_chunks_empty")
    return chunks


def document_content_type(file_name: str, content_type: str | None) -> str:
    extension = _document_extension(file_name)
    if extension == ".csv":
        return "text/csv"
    if extension == ".md":
        return "text/markdown"
    if extension == ".txt":
        return "text/plain"
    if extension == ".pdf":
        return "application/pdf"
    if extension == ".docx":
        return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    if extension == ".xlsx":
        return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    if extension == ".pptx":
        return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    return str(content_type or "application/octet-stream")


def safe_document_file_name(file_name: str | None) -> str:
    normalized = str(file_name or "").replace("\\", "/")
    safe_name = Path(normalized).name.strip()
    if not safe_name or safe_name in {".", ".."}:
        raise DocumentIngestionError("document_file_name_required")
    return safe_name


def _document_extension(file_name: str) -> str:
    return Path(file_name).suffix.lower()


def _decode_text(raw_bytes: bytes) -> str:
    try:
        return raw_bytes.decode("utf-8-sig")
    except UnicodeDecodeError:
        try:
            return raw_bytes.decode("cp932")
        except UnicodeDecodeError as error:
            raise DocumentIngestionError("document_text_encoding_unsupported") from error


def _extract_docx_text(raw_bytes: bytes) -> str:
    document = DocxDocument(BytesIO(raw_bytes))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if values:
                parts.append(" / ".join(values))
    return "\n".join(parts)


def _extract_xlsx_text(raw_bytes: bytes) -> str:
    workbook = load_workbook(BytesIO(raw_bytes), read_only=True, data_only=True)
    try:
        parts: list[str] = []
        for worksheet in workbook.worksheets:
            for row in worksheet.iter_rows(values_only=True):
                values = [str(value).strip() for value in row if value is not None and str(value).strip()]
                if values:
                    parts.append(" / ".join(values))
        return "\n".join(parts)
    finally:
        workbook.close()


def _extract_pptx_text(raw_bytes: bytes) -> str:
    presentation = Presentation(BytesIO(raw_bytes))
    parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text.strip()
            if text:
                parts.append(text)
    return "\n".join(parts)


def _update_document(document_id: str, **changes: Any) -> dict[str, Any]:
    document = store.get("documents", document_id)
    if not document:
        raise DocumentIngestionError("document_not_found")
    document.update(changes)
    document["updatedAt"] = utc_now()
    store.upsert("documents", document)
    return document
